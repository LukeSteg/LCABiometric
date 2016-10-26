#!/usr/bin/python

import xlrd
from shutil import copyfile
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches,Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION

def parseSlide(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            frame = shape.text_frame
            print(frame.text)
            text = frame.text
            text = text.replace('#VAR:PARTICIPANT_COUNT','48')
            text = text.replace('#VAR:PERCENT_FEMALE','67')
            frame.text = text
            #shape = frame
            print(frame.text)
            print(shape.text_frame.text)
            if (shape.text_frame.text == '#CHART:PREFAB1'):
                print(shape.text_frame.text)
                shape.text_frame.clear
                x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
                slide.shapes.add_chart( XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)


outputFileName = 'exampleOutput.pptx'
inputTemplateFileName = 'exampleTemplate3.pptx'

copyfile(inputTemplateFileName, outputFileName)

#outputFile = open(outputFileName)
#prs = Presentation(outputFile)
prs = Presentation(inputTemplateFileName)

slide = prs.slides[0]

for x in range(0, 6):
    slide = prs.slides[x]
    print(len(list(slide.placeholders._element.iter_ph_elms())))
    for shape in slide.placeholders:
        print(shape.placeholder_format.idx, shape.name)


slide = prs.slides[1]
parseSlide(slide)


book = xlrd.open_workbook("ExampleDataset.xlsx")
aggregatedSheet = book.sheet_by_index(25)
chart_data = ChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

peopleAtGoal = 0
peopleNotAtGoal = 0
numberOfPeople = 0
for i in range(aggregatedSheet.nrows - 1):
    numberOfPeople += 1
    if( aggregatedSheet.cell_value(rowx = i + 1, colx = 5) == "Goal"):
        peopleAtGoal += 1;
        print("here" , peopleAtGoal , i );
    else:
        peopleNotAtGoal += 1
        percentageAtGoal = 100*(float(peopleAtGoal)/numberOfPeople)
        percentageNotAtGoal = 100*(float(peopleNotAtGoal)/numberOfPeople)
        print(percentageAtGoal,percentageNotAtGoal)
        chart_data = ChartData()
        chart_data.categories = ['Goal','Less than Goal']
        chart_data.add_series('Series 1',(percentageAtGoal, percentageNotAtGoal))

slide = prs.slides[2]
parseSlide(slide)        

prs.save(outputFileName)



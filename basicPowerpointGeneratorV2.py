#!/usr/bin/python

import xlrd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches,Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION

def AddProteinNutritionSlide(prs,goal, less_than_goal, greater_than_goal):
	slide = prs.slides.add_slide(prs.slide_layouts[3]) #title then 2 body layout
	slide.shapes.title.text = "Meeting Daily Nutrition Requirements"


	shapes = slide.shapes
	bodyshape1 = shapes.placeholders[1]

	x, y, cx, cy = Inches(1), Inches(1), Inches(5), Inches(4)

	chart_data = ChartData()
	chart_data.categories = ['At Recommended Amount', 'Less Than Recommended Amount','Greater Than Recommended Amount']
	chart_data.add_series('Daily Intake of Protein', (goal,less_than_goal,greater_than_goal))

	chart = slide.shapes.add_chart(
	    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
	).chart

	chart.has_legend = True
	chart.legend.position = XL_LEGEND_POSITION.BOTTOM
	chart.legend.include_in_layout = False

	chart.plots[0].has_data_labels = True
	data_labels = chart.plots[0].data_labels
	data_labels.number_format = '0%'
	data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

	
	left = top = width = height = Inches(1)
	bodyshape2 = shapes.placeholders[2]
	tf = bodyshape2.text_frame

	p = tf.add_paragraph()
	p.font.size = Pt(20)
	p.text = "Inadequate intake of nutrient-dense foods can lead to nutrient deficiencies, impairs worker productivity, and contributes to disease risk.  Even small positive dietary changes can have a profound effect on overall health and wellbeing."

	return slide


def AddDailyNutritionSlide(prs,percent_fruit, percent_grain, percent_vegetable, percent_calcium):
	slide = prs.slides.add_slide(prs.slide_layouts[5]) #title only layout
	slide.shapes.title.text = "Meeting Daily Nutrition Requirements"


	#txBox = slide.shapes.add_textbox(left, top, width, height)
	shapes = slide.shapes

	#TODO Compress production of charts here
	################################################
	x, y, cx, cy = Inches(1), Inches(2), Inches(3), Inches(3)

	chart_data = ChartData()
	chart_data.categories = ['Goal', 'Less Than Goal']
	chart_data.add_series('Daily Intake of Fruits', (percent_fruit,1-percent_fruit))

	chart = slide.shapes.add_chart(
	    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
	).chart

	chart.has_legend = True
	chart.legend.position = XL_LEGEND_POSITION.BOTTOM
	chart.legend.include_in_layout = False

	chart.plots[0].has_data_labels = True
	data_labels = chart.plots[0].data_labels
	data_labels.number_format = '0%'
	data_labels.position = XL_LABEL_POSITION.OUTSIDE_END


	################################################

	x, y, cx, cy = Inches(5), Inches(2), Inches(4), Inches(3)

	chart_data = ChartData()
	chart_data.categories = ['Goal', 'Less Than Goal']
	chart_data.add_series('Daily Intake of Whole Grain', (percent_grain,1-percent_grain))

	chart = slide.shapes.add_chart(
	    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
	).chart

	chart.has_legend = True
	chart.legend.position = XL_LEGEND_POSITION.BOTTOM
	chart.legend.include_in_layout = False

	chart.plots[0].has_data_labels = True
	data_labels = chart.plots[0].data_labels
	data_labels.number_format = '0%'
	data_labels.position = XL_LABEL_POSITION.OUTSIDE_END


	################################################

	x, y, cx, cy = Inches(0.5), Inches(5), Inches(4), Inches(3)

	chart_data = ChartData()
	chart_data.categories = ['Goal', 'Less Than Goal']
	chart_data.add_series('Daily Intake of Vegetables', (percent_vegetable,1-percent_vegetable))

	chart = slide.shapes.add_chart(
	    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
	).chart

	chart.has_legend = True
	chart.legend.position = XL_LEGEND_POSITION.BOTTOM
	chart.legend.include_in_layout = False

	chart.plots[0].has_data_labels = True
	data_labels = chart.plots[0].data_labels
	data_labels.number_format = '0%'
	data_labels.position = XL_LABEL_POSITION.OUTSIDE_END


	################################################

	x, y, cx, cy = Inches(5), Inches(5), Inches(4), Inches(3)

	chart_data = ChartData()
	chart_data.categories = ['Goal', 'Less Than Goal']
	chart_data.add_series('Daily Intake of Calcium', (percent_calcium,1-percent_calcium))

	chart = slide.shapes.add_chart(
	    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
	).chart

	chart.has_legend = True
	chart.legend.position = XL_LEGEND_POSITION.BOTTOM
	chart.legend.include_in_layout = False

	chart.plots[0].has_data_labels = True
	data_labels = chart.plots[0].data_labels
	data_labels.number_format = '0%'
	data_labels.position = XL_LABEL_POSITION.OUTSIDE_END





	return slide 


def AddParticipationSlide(prs,num_participants, percent_fem, average_age):
	slide = prs.slides.add_slide(prs.slide_layouts[1]) #bullet layout
	slide.shapes.title.text = "Who Participated?"


	left = top = width = height = Inches(1)
	#txBox = slide.shapes.add_textbox(left, top, width, height)
	shapes = slide.shapes
	bodyshape = shapes.placeholders[1]
	tf = bodyshape.text_frame

	line1 = "{} Individuals completed fasting biometric screening".format(num_participants)
	line2 = "Female Participants: {}%".format(percent_fem)
	line3 = "Male Participants: {}%".format(100-percent_fem)
	line4 = "Average Age: {} years".format(average_age)

	#TODO make this a loop
	p = tf.add_paragraph()
	p.font.size = Pt(40)
	p.text = line1

	p = tf.add_paragraph()
	p.font.size = Pt(40)
	p.text = line2

	p = tf.add_paragraph()
	p.font.size = Pt(40)
	p.text = line3

	p = tf.add_paragraph()
	p.font.size = Pt(40)
	p.text = line4

	return slide 


#Reading a book
book = xlrd.open_workbook("ExampleDataset.xlsx")

#Using a book
print("The number of worksheets is {0}".format(book.nsheets))
print("Worksheet name(s): {0}".format(book.sheet_names()))

#Using a sheet
print("----------- Sheet 1 -----------")
sh = book.sheet_by_index(0)
print("{0} {1} {2}".format(sh.name, sh.nrows, sh.ncols))
print("Cell A4 is {0}".format(sh.cell_value(rowx=4, colx=0)))
for rx in range(sh.nrows):
    print(sh.row(rx))
print("--------------------------------")


#Iterating all rows of all sheets
print("------- All Sheets -------")
for ns in range(book.nsheets):
	print "----------- Sheet",ns,"-----------"
	sh = book.sheet_by_index(ns)
	print("{0} {1} {2}".format(sh.name, sh.nrows, sh.ncols))
	for rx in range(sh.nrows):
	    print(sh.row(rx))
	print("--------------------------------")

aggregatedSheet = book.sheet_by_index(25)


# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = ChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

peopleAtGoal = 0
peopleNotAtGoal = 0
numberOfPeople = aggregatedSheet.nrows - 1 
for i in range(aggregatedSheet.nrows - 1):
    numberOfPeople += 1
    if( aggregatedSheet.cell_value(rowx = i + 1, colx = 5) == "Goal"):
        peopleAtGoal += 1;
        print("here" , peopleAtGoal , i );
    else:
        peopleNotAtGoal += 1

percentageAtGoal = 100*(float(peopleAtGoal)/numberOfPeople)
percentageNotAtGoal = 100*(float(peopleNotAtGoal)/numberOfPeople)
print(percentageAtGoal,percentageNotAtGoal)
chart_data = ChartData()
chart_data.categories = ['Goal','Less than Goal']
chart_data.add_series('Series 1',(percentageAtGoal, percentageNotAtGoal))

# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )

AddParticipationSlide(prs,48,67,47)
AddDailyNutritionSlide(prs,0.69,0.51,0.63,0.43)
AddProteinNutritionSlide(prs,0.78,0.18,0.04)
prs.save('chart-01.pptx')

def getParticipationSlide(presentation, numberOfParticipants, percentageFemale, averageAge)
    participationSlide = presentation.slide_layouts[1]
    titleText = "Who Participated?"
    titleTextFrame.text = titleText

    contentTextFrame.text = """%i Individuals completed fasting biometric screening\nFemale Participants: %i\%\nMale
    Participants: %i\%\nAverage Age: %i years""" %(numberOfParticipants, percentageFemale, 100 - percentageFemale,
    averageAge)

    

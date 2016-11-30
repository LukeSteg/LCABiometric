#!/usr/bin/python

import xlrd
from shutil import copyfile
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches,Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION
from Parser import parse
from Parser import parseTable

def parseSlide(slide,fileDict):
    for shape in slide.shapes:
        if shape.has_text_frame:
            parse(slide,shape,shape,fileDict)
        if shape.has_table:
            parseTable(slide, shape, shape, fileDict)
            
def createOutput(outputFileName, inputTemplateFileName, DataSheetDict):

	prs = Presentation(inputTemplateFileName)

	for slide in prs.slides:
		parseSlide(slide, DataSheetDict)
	
	prs.save(outputFileName)



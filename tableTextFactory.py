#!/usr/bin/python

import xlrd
from textFactory import textFactory
from genericFactory import AGGREGATE_SHEET_NAME
from pptx.chart.data import ChartData

class tableTextFactory(textFactory):

    def __init__(self, slideRef, shapeRef, cellRef):
        super(tableTextFactory,self).__init__(slideRef, shapeRef)
        self.cellRef = cellRef
        print 'cell contents' + cellRef.text_frame.text

    def generateShape(self):
        print self.outputText 
        queryStartIndex = self.cellRef.text_frame.text.index('#{')
        queryEndIndex = self.cellRef.text_frame.text.index('}')
        self.cellRef.text = self.cellRef.text_frame.text[0:queryStartIndex] + self.outputText + self.cellRef.text_frame.text[queryEndIndex+1: len(self.cellRef.text_frame.text)]


from pptx.util import Inches
import pieChartFactory
import barChartFactory
import textFactory
import lineChartFactory
import tableTextFactory

def get_arg_dict(tokens):
    arg_dict = {}
    for token in tokens:
        tkn_type = token.split(':')[0].strip().upper()
        tkn_value = token.split(':')[1].strip().upper()
        arg_dict[tkn_type] = tkn_value
    return arg_dict

def generate_pie_chart(slide,shape,tokens,fileDict):
    pf = pieChartFactory.pieChartFactory(slide,shape)
    
    arg_dict = get_arg_dict(tokens)

    defaultFactoryActions(pf, arg_dict, fileDict);    
#    if('TITLE' in arg_dict):
#        pf.setTitle(arg_dict['TITLE'])
    if('COLUMN' in arg_dict):
        pf.setColumn(arg_dict['COLUMN'], fileDict)
        pf.getDataFromColumn(pf.getFileFromDict(fileDict))

    return pf.generateShape()

def generate_bar_chart(slide,shape,tokens,fileDict):
    bf = barChartFactory.barChartFactory(slide,shape)
    
    arg_dict = get_arg_dict(tokens)

    defaultFactoryActions(bf, arg_dict, fileDict);    
    if('COLUMN' in arg_dict):
        bf.setColumn(arg_dict['COLUMN'], fileDict)
        bf.getDataFromColumn(bf.getFileFromDict(fileDict))

    return bf.generateShape()

def generate_line_chart(slide,shape,tokens,fileDict):
    lf = lineChartFactory.lineChartFactory(slide,shape)

    arg_dict = get_arg_dict(tokens)

    #MUST GO BEFORE 'COLUMN'
    if('SURVEYCOUNT' in arg_dict):
        lf.setNumberOfBooks(int(arg_dict['SURVEYCOUNT']))
    
    defaultFactoryActions(lf, arg_dict, fileDict);   
    if('COLUMN' in arg_dict):
        lf.setColumnName(arg_dict['COLUMN'])
        lf.getDataFromColumn(fileDict)

    return lf.generateShape()

def generate_text(slide,shape,tokens,fileDict):
    tf = textFactory.textFactory(slide,shape)

    arg_dict = get_arg_dict(tokens)
        
    defaultFactoryActions(tf, arg_dict, fileDict);    
    if('COUNTORPERCENT' in arg_dict):
        tf.setOutputVarType(arg_dict['COUNTORPERCENT'])
    if('COLUMN' in arg_dict):
        tf.setColumn(arg_dict['COLUMN'], fileDict)
        tf.getDataFromColumn(tf.getFileFromDict(fileDict))
    if('VARIABLE' in arg_dict):
        tf.computeOutputVar(arg_dict['VARIABLE'])
    if('FONT' in arg_dict):
        tf.setTextSize(int(arg_dict['FONT']))
   
    return tf.generateShape()

def generate_table_text(slide, shape, tokens, fileDict, cellRef):
    tf = tableTextFactory.tableTextFactory(slide, shape, cellRef)

    arg_dict = get_arg_dict(tokens)

    if('BOOK' in arg_dict):
    	tf.setBook(int(arg_dict['BOOK']))
    if('COLUMN' in arg_dict):
        tf.setColumn(arg_dict['COLUMN'], fileDict)
        tf.getDataFromColumn(tf.getFileFromDict(fileDict))
    if('VARIABLE' in arg_dict):
        tf.computeOutputVar(arg_dict['VARIABLE'])
    if('FONT' in arg_dict):
        tf.setTextSize(int(arg_dict['FONT']))
   
    return tf.generateShape()



def parse(slide,shape0,shape,fileDict):
    frame = shape.text_frame
    
    text = frame.text.strip().encode('utf8')
    print "parsing %s" % text  
    #define if while actually worked *cough* zane *cough*  
    while containsQueryString(text):
        text = getQueryString(text)
        print 'Query String ' + text
        tokens = text.split(',')
        print type(tokens[0])
        tokens = map(str,tokens)
        print type(tokens[0])
        tokens = map(str.strip,tokens)
        tokens = map(str.upper,tokens)
        
        fig_type = tokens[0].split(':')[1]

        switch = {'PIE CHART':generate_pie_chart , 'BAR CHART':generate_bar_chart , 'LINE CHART':generate_line_chart , 'TEXT':generate_text }
         
        new_shape = switch[fig_type](slide,shape,tokens,fileDict)
        #update for while loop
        text = frame.text.strip()

def parseTable(slide, shape0, shape, fileDict):
    table = shape.table
    
    print 'parsing table'
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            cell = table.cell(r,c)
            cellText = cell.text_frame.text
            #define if while take 2
            while(containsQueryString(cellText)):
                text = getQueryString(cellText)
                print 'Query String ' + text
                tokens = text.split(',');
                print type(tokens[0]);
                tokens = map(str,tokens)
                print type(tokens[0]);
                tokens = map(str.strip,tokens);
                tokens = map(str.upper,tokens);
        
                generate_table_text(slide, shape, tokens, fileDict, cell)
                #update for while loop that might work
                cellText = cell.text_frame.text

def containsQueryString(text):
    result = False
    if(('#{' in text) and ('}' in text)):
        if(text.index('#{') < text.index('}')):
            result = True
    return result

def getQueryString(text):
    #modify to create multiple query strings?
    if(('#{' in text) and ('}' in text)):
        startIndex = text.index('#{')
        endIndex = text.index('}')
        return text[startIndex+2:endIndex]
    else:
        print 'WARNING query string not found when possibly expected'
        return ''

def defaultFactoryActions(factory, arg_dict, fileDict):

    if('X' in arg_dict):
        factory.setX(arg_dict['X'])
    if('Y' in arg_dict):
        factory.setY(arg_dict['Y'])
    if('CX' in arg_dict):
        factory.setCX(arg_dict['CX'])
    if('CY' in arg_dict):
        factory.setCY(arg_dict['CY'])
    if('BOOK' in arg_dict):
    	factory.setBook(int(arg_dict['BOOK']))
 
